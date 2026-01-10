import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

def read_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        reader = PdfReader(file_path)
        return [page.extract_text() for page in reader.pages]
    
    elif ext == '.docx':
        doc = Document(file_path)
        return [p.text for p in doc.paragraphs if p.text.strip()]
    
    elif ext == '.pptx':
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            text_runs.append(slide_text.strip())
        return text_runs

    elif ext in ['.png', '.jpg', '.jpeg']:
        # Basic OCR for images
        text = pytesseract.image_to_string(Image.open(file_path))
        return [text]
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")